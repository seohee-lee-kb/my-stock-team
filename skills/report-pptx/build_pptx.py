# -*- coding: utf-8 -*-
"""
report-pptx 빌더
================
reports/{종목명}.md (종목 리서치) → reports/{종목명}.pptx (디자인된 리포트)

사용:
    python build_pptx.py <input.md> [output.pptx]

설계 원칙 (CLAUDE.md / 스킬 사양 준수)
- 슬라이드 순서 고정: 표지 → 종목 개요 → 재무 요약 → 가격/추세(차트)
  → 뉴스·심리 → 리스크 → 한 줄 종합
- 포인트색 KB 옐로우(#FFBC00) + 본문 그레이/화이트, 차분한 금융 리포트 톤
- 한글 폰트는 '맑은 고딕' 하나로 고정 (latin/ea/cs 모두 지정 → 글자 깨짐 방지)
- 모든 수치 옆 (출처/기준일) 표기. 출처 라인을 캡션으로 추출해 슬라이드 하단에 표기
- 매수/매도·목표가 단정 표현은 스크럽으로 제거
- 표가 슬라이드 밖으로 넘치지 않도록 행 수 상한 + 넘치면 '이하 생략' 처리
"""
from __future__ import annotations

import datetime as _dt
import os
import re
import sys
import tempfile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# 디자인 상수
# ---------------------------------------------------------------------------
KB_YELLOW = RGBColor(0xFF, 0xBC, 0x00)   # 포인트색
INK = RGBColor(0x33, 0x33, 0x33)         # 본문 진회색
SUBINK = RGBColor(0x70, 0x70, 0x70)      # 보조 회색 (캡션/출처)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TABLE_HEAD_BG = RGBColor(0x33, 0x33, 0x33)   # 표 헤더 진회색
TABLE_ALT_BG = RGBColor(0xF4, 0xF4, 0xF4)    # 표 짝수행 옅은 회색
TABLE_LINE = RGBColor(0xD9, 0xD9, 0xD9)

FONT = "맑은 고딕"
# matplotlib은 한글 폰트명 해석에 실패할 수 있어 파일 경로로 직접 등록한다.
def _find_korean_font():
    candidates = [
        r"C:\Windows\Fonts\malgun.ttf",                      # Windows 맑은 고딕
        "/Library/Fonts/AppleGothic.ttf",                      # macOS
        "/System/Library/Fonts/Supplemental/AppleGothic.ttf",  # macOS
        "/usr/share/fonts/truetype/nanum/NanumGothic.ttf",     # Linux 나눔고딕
        "/usr/share/fonts/opentype/noto/NotoSansCJKkr-Regular.otf",  # Linux Noto
    ]
    for _p in candidates:
        if os.path.exists(_p):
            return _p
    return ""
MALGUN_TTF = _find_korean_font()  # 없으면 "" → 차트/PDF는 기본 폰트로 폴백

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.7)
CONTENT_TOP = Inches(1.55)               # 제목 밴드 아래 본문 시작
CONTENT_W = SLIDE_W - 2 * MARGIN

MAX_TABLE_ROWS = 8                        # 헤더 제외 최대 데이터 행
MAX_BULLETS = 7                           # 슬라이드당 최대 불릿

# 매수/매도·목표가 단정 표현 → 제거 (출처 attribution 없는 자체 단정만 겨냥)
GUARDRAIL_PATTERNS = [
    re.compile(r"(매수|매도|비중확대|비중축소)\s*(추천|권고|의견|하(세요|십시오|시기))"),
    re.compile(r"목표[가주]\s*(는|:)?\s*[\d,]+\s*원"),
    re.compile(r"(강력|적극)\s*(매수|매도)"),
]


# ---------------------------------------------------------------------------
# 폰트 유틸 — latin/ea/cs 모두 '맑은 고딕'으로 강제 (한글 깨짐 방지)
# ---------------------------------------------------------------------------
def style_run(run, size, bold=False, color=INK, font=FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", font)


def add_textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tb, tf


def write_para(tf, text, size, bold=False, color=INK, first=False,
               align=PP_ALIGN.LEFT, bullet=False, space_after=6):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    if bullet:
        run = p.add_run()
        run.text = "•  "
        style_run(run, size, bold=True, color=KB_YELLOW)
    run = p.add_run()
    run.text = text
    style_run(run, size, bold=bold, color=color)
    return p


# ---------------------------------------------------------------------------
# 공통 슬라이드 뼈대
# ---------------------------------------------------------------------------
def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_rect(slide, left, top, width, height, color):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_title_band(slide, title, idx=None, total=None):
    """제목 + KB 옐로우 언더라인 악센트 + 우상단 페이지 표기."""
    # 좌측 옐로우 세로 악센트
    add_rect(slide, MARGIN, Inches(0.55), Inches(0.12), Inches(0.55), KB_YELLOW)
    _, tf = add_textbox(slide, MARGIN + Inches(0.28), Inches(0.5),
                        CONTENT_W - Inches(0.28), Inches(0.7))
    write_para(tf, title, 26, bold=True, color=INK, first=True)
    # 제목 하단 옐로우 라인
    add_rect(slide, MARGIN, Inches(1.28), CONTENT_W, Pt(2.2), KB_YELLOW)
    if idx is not None and total is not None:
        _, ptf = add_textbox(slide, SLIDE_W - Inches(2.0), Inches(0.6),
                             Inches(1.3), Inches(0.4))
        write_para(ptf, f"{idx} / {total}", 11, color=SUBINK,
                   first=True, align=PP_ALIGN.RIGHT)


def add_footer(slide, source_text):
    """하단 출처/기준일 캡션 + '학습용' 고지."""
    _, tf = add_textbox(slide, MARGIN, SLIDE_H - Inches(0.55),
                        CONTENT_W, Inches(0.45))
    note = "본 자료는 학습용 분석이며 투자 자문이 아닙니다."
    txt = f"{source_text}   ·   {note}" if source_text else note
    write_para(tf, txt, 8.5, color=SUBINK, first=True)


# ---------------------------------------------------------------------------
# 마크다운 파싱
# ---------------------------------------------------------------------------
def scrub_guardrail(text: str) -> str:
    for pat in GUARDRAIL_PATTERNS:
        text = pat.sub("[판단 근거는 본문 참조]", text)
    return text


def split_sections(md: str):
    """## / # 헤딩 기준으로 (title, body) 섹션 리스트로 분해."""
    lines = md.splitlines()
    sections, cur_title, cur_body = [], None, []
    h1 = None
    for ln in lines:
        m = re.match(r"^(#{1,3})\s+(.*)$", ln)
        if m:
            level, title = len(m.group(1)), m.group(2).strip()
            if level == 1 and h1 is None:
                h1 = re.sub(r"\s*리서치\s*$", "", title).strip()
            if cur_title is not None or cur_body:
                sections.append((cur_title, "\n".join(cur_body).strip()))
            cur_title, cur_body = title, []
        else:
            cur_body.append(ln)
    if cur_title is not None or cur_body:
        sections.append((cur_title, "\n".join(cur_body).strip()))
    return h1, sections


def parse_tables(text: str):
    """본문에서 마크다운 파이프 표를 모두 추출 → [ [row,...], ... ]."""
    tables, cur = [], []
    for ln in text.splitlines():
        s = ln.strip()
        if s.startswith("|") and s.endswith("|"):
            cells = [c.strip() for c in s.strip("|").split("|")]
            if re.match(r"^[\s:\-|]+$", s.replace("|", "")):  # 구분선
                continue
            cur.append(cells)
        else:
            if cur:
                tables.append(cur)
                cur = []
    if cur:
        tables.append(cur)
    return [t for t in tables if len(t) >= 2]


def extract_sources(text: str) -> str:
    """(출처...) / *(출처...)* / 기준일 캡션을 모아 슬라이드 하단 문구로."""
    hits = []
    for m in re.finditer(r"\(?\*?\(출처[^)]*\)\*?\)?", text):
        hits.append(re.sub(r"[*]", "", m.group(0)).strip())
    # 중복 제거 + 순서 유지
    seen, out = set(), []
    for h in hits:
        h = h.strip("()").strip()
        if h and h not in seen:
            seen.add(h)
            out.append(h)
    return " / ".join(out[:3])


def clean_bullets(text: str):
    """표·출처캡션·구분선 제거 후, 불릿/문장 리스트 반환."""
    out = []
    for raw in text.splitlines():
        s = raw.strip()
        if not s or s.startswith("|") or s.startswith("---"):
            continue
        if re.match(r"^\*?\(?출처", s) or s.startswith("*(출처"):
            continue
        s = re.sub(r"^[-*•]\s*", "", s)          # 불릿 마커 제거
        s = re.sub(r"^\*\*(.+?)\*\*", r"\1", s)   # 굵게 마커 정리
        s = s.replace("**", "").replace("`", "")
        s = s.replace("*(", "(").replace(")*", ")").replace("*", "")  # 인라인 강조/캡션 별표 정리
        s = scrub_guardrail(s).strip()
        if s:
            out.append(s)
    return out


# ---------------------------------------------------------------------------
# 섹션 분류 (제목 키워드 → 슬라이드 버킷)
# ---------------------------------------------------------------------------
def classify(title: str) -> str | None:
    if not title:
        return None
    t = title.lower()
    rules = [
        ("overview", ["개요", "프로필", "overview", "소개"]),
        ("finance", ["재무", "실적", "financ", "밸류"]),
        ("price", ["가격", "추세", "차트", "주가", "기술", "price", "chart"]),
        ("news", ["뉴스", "심리", "이슈", "sentiment", "news"]),
        ("risk", ["리스크", "위험", "risk"]),
        ("summary", ["종합", "한 줄", "한줄", "결론", "의견", "summary"]),
        ("cover", ["표지", "cover"]),
    ]
    for bucket, kws in rules:
        if any(k in t for k in kws):
            return bucket
    return None


def collect(sections):
    """버킷별로 body 텍스트를 합쳐 dict 반환."""
    buckets = {}
    for title, body in sections:
        b = classify(title)
        if b is None:
            continue
        buckets.setdefault(b, []).append((title, body))
    return buckets


def find_meta(sections):
    """기준일/작성일/작성 관점 메타 추출."""
    joined = "\n".join(b for _, b in sections)
    date = None
    m = re.search(r"(기준일|작성일)\s*[:：]?\s*([0-9]{4}-[0-9]{2}-[0-9]{2})", joined)
    if m:
        date = m.group(2)
    perspective = None
    m = re.search(r"작성\s*관점\s*[:：]?\s*([^\n|]+)", joined)
    if m:
        perspective = m.group(1).strip().strip("|").strip()
    return date, perspective


# ---------------------------------------------------------------------------
# 표 렌더링 (자동 폭/행 제한)
# ---------------------------------------------------------------------------
def render_table(slide, rows, top, max_rows=MAX_TABLE_ROWS):
    if not rows:
        return top
    header, data = rows[0], rows[1:]
    truncated = False
    if len(data) > max_rows:
        data = data[:max_rows]
        truncated = True
    body = [header] + data + ([["…", *([""] * (len(header) - 1))]] if truncated else [])

    n_rows, n_cols = len(body), len(header)
    # 폭: 첫 열 넓게, 나머지 균등
    first_w = Inches(2.6) if n_cols > 1 else CONTENT_W
    rest_w = (CONTENT_W - first_w) / max(1, n_cols - 1) if n_cols > 1 else 0
    height = Inches(min(0.45 * n_rows + 0.1, 4.8))

    gtable = slide.shapes.add_table(n_rows, n_cols, MARGIN, top, CONTENT_W, height)
    table = gtable.table
    # python-pptx 기본 스타일 밴딩 끄기
    table.first_row = False
    table.horz_banding = False

    for c in range(n_cols):
        table.columns[c].width = int(first_w if c == 0 else rest_w)

    # 셀 폰트 크기: 열 수 많으면 축소
    fsize = 12 if n_cols <= 4 else (10.5 if n_cols <= 6 else 9)

    for r in range(n_rows):
        for c in range(n_cols):
            cell = table.cell(r, c)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            txt = body[r][c] if c < len(body[r]) else ""
            txt = txt.replace("**", "")
            tf = cell.text_frame
            tf.word_wrap = True
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            run = p.add_run()
            run.text = txt
            if r == 0:  # 헤더
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEAD_BG
                style_run(run, fsize, bold=True, color=WHITE)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG if r % 2 == 0 else WHITE
                style_run(run, fsize, bold=(c == 0), color=INK)
    return top + height + Inches(0.15)


# ---------------------------------------------------------------------------
# 차트 렌더링 (matplotlib → PNG → 슬라이드)
# ---------------------------------------------------------------------------
PERIOD_ORDER = ["1주", "1개월", "1달", "3개월", "6개월", "1년", "연초"]


def _num(s):
    m = re.search(r"[-+]?\d[\d,]*\.?\d*", s.replace(",", ""))
    return float(m.group(0)) if m else None


def build_chart(price_bodies, out_dir):
    """가격/추세 섹션 표에서 기간별 변동률(%)을 찾아 막대차트 PNG 생성.
    못 찾으면 임의 2열 숫자표를 막대로. 그래도 없으면 None 반환."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    from matplotlib import font_manager

    # 한글 폰트를 파일 경로로 등록(폰트명 해석 실패/글자 깨짐 방지)
    mpl_font = FONT
    if os.path.exists(MALGUN_TTF):
        font_manager.fontManager.addfont(MALGUN_TTF)
        mpl_font = font_manager.FontProperties(fname=MALGUN_TTF).get_name()
    plt.rcParams["font.family"] = mpl_font
    plt.rcParams["axes.unicode_minus"] = False

    labels, values, title = [], [], "기간별 변동률"
    for _, body in price_bodies:
        for tbl in parse_tables(body):
            for row in tbl[1:]:
                if len(row) < 2:
                    continue
                lab = row[0]
                if any(k in lab for k in PERIOD_ORDER) and "변동" in lab or \
                   re.search(r"(1주|1개월|3개월|6개월|1년|연초).*변동", lab):
                    v = _num(row[1])
                    if v is not None:
                        labels.append(lab.replace(" 변동률", "").replace("변동률", "").strip())
                        values.append(v)
        if labels:
            break

    # 폴백: 기간 변동률 표가 없으면 첫 숫자형 2열 표 사용
    if not labels:
        for _, body in price_bodies:
            for tbl in parse_tables(body):
                for row in tbl[1:]:
                    if len(row) >= 2 and _num(row[1]) is not None and not re.search(r"\d{4,}", row[0]):
                        labels.append(row[0][:8])
                        values.append(_num(row[1]))
                if labels:
                    title = (tbl[0][1] if len(tbl[0]) > 1 else "지표")
                    break
            if labels:
                break
    if not labels:
        return None

    labels, values = labels[:8], values[:8]
    fig, ax = plt.subplots(figsize=(8.6, 3.9), dpi=150)
    colors = ["#FFBC00" if v >= 0 else "#9A9A9A" for v in values]
    bars = ax.bar(labels, values, color=colors, edgecolor="white", width=0.62)
    ax.axhline(0, color="#666666", linewidth=0.8)
    ax.set_title(title, fontsize=14, fontweight="bold", color="#333333", pad=12)
    ax.spines[["top", "right"]].set_visible(False)
    ax.spines[["left", "bottom"]].set_color("#CCCCCC")
    ax.tick_params(axis="x", colors="#555555", labelsize=11, pad=14)
    ax.tick_params(axis="y", colors="#555555", labelsize=11)
    ax.grid(axis="y", color="#EEEEEE", linewidth=0.8)
    # 음수 막대 라벨이 x축 카테고리 라벨과 겹치지 않도록 y범위에 여유를 둔다
    vmin, vmax = min(values + [0]), max(values + [0])
    pad = (vmax - vmin) * 0.12 or 1
    ax.set_ylim(vmin - pad, vmax + pad)
    for b, v in zip(bars, values):
        ax.annotate(f"{v:+.1f}", (b.get_x() + b.get_width() / 2, v),
                    ha="center", va="bottom" if v >= 0 else "top",
                    fontsize=10, color="#333333",
                    xytext=(0, 4 if v >= 0 else -4), textcoords="offset points")
    fig.tight_layout()
    path = os.path.join(out_dir, "_chart_tmp.png")
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return path


# ---------------------------------------------------------------------------
# 슬라이드 빌더들
# ---------------------------------------------------------------------------
def slide_cover(prs, name, date, perspective):
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    # 좌측 옐로우 대형 악센트 바
    add_rect(s, 0, 0, Inches(0.45), SLIDE_H, KB_YELLOW)
    add_rect(s, MARGIN, Inches(2.55), Inches(1.6), Pt(5), KB_YELLOW)
    _, tf = add_textbox(s, MARGIN, Inches(2.75), SLIDE_W - 2 * MARGIN, Inches(1.6))
    write_para(tf, name, 46, bold=True, color=INK, first=True)
    write_para(tf, "종목 리서치 리포트", 20, color=SUBINK, space_after=4)
    _, tf2 = add_textbox(s, MARGIN, Inches(5.1), SLIDE_W - 2 * MARGIN, Inches(1.2))
    if date:
        write_para(tf2, f"작성일 / 기준일 : {date}", 14, color=INK, first=True)
    if perspective:
        write_para(tf2, f"작성 관점 : {perspective}", 14, color=SUBINK)
    add_footer(s, "")
    return s


def slide_text(prs, title, bodies, idx, total, as_bullets=True):
    s = blank_slide(prs)
    add_title_band(s, title, idx, total)
    src = " / ".join(filter(None, (extract_sources(b) for _, b in bodies)))
    lines = []
    for _, body in bodies:
        lines.extend(clean_bullets(body))
    _, tf = add_textbox(s, MARGIN, CONTENT_TOP, CONTENT_W, SLIDE_H - CONTENT_TOP - Inches(0.7))
    shown = lines[:MAX_BULLETS]
    for i, ln in enumerate(shown):
        # 너무 긴 줄은 잘라 표시
        ln = ln if len(ln) <= 95 else ln[:93] + "…"
        write_para(tf, ln, 14, color=INK, first=(i == 0), bullet=as_bullets, space_after=10)
    add_footer(s, src)
    return s


def slide_finance(prs, bodies, idx, total):
    s = blank_slide(prs)
    add_title_band(s, "재무 요약 (최근 3개년)", idx, total)
    src = " / ".join(filter(None, (extract_sources(b) for _, b in bodies)))
    top = CONTENT_TOP
    # 가장 큰 표 1개 우선 렌더
    tables = []
    for _, body in bodies:
        tables.extend(parse_tables(body))
    if tables:
        tables.sort(key=lambda t: (len(t), len(t[0])), reverse=True)
        top = render_table(s, tables[0], top)
    # 표 아래 핵심 코멘트 1~2줄
    comments = []
    for _, body in bodies:
        for ln in clean_bullets(body):
            if "코멘트" in ln or len(ln) > 25:
                comments.append(ln)
    if comments and top < SLIDE_H - Inches(1.4):
        _, tf = add_textbox(s, MARGIN, top, CONTENT_W, SLIDE_H - top - Inches(0.7))
        for i, ln in enumerate(comments[:2]):
            ln = ln if len(ln) <= 95 else ln[:93] + "…"
            write_para(tf, ln, 12.5, color=INK, first=(i == 0), bullet=True, space_after=8)
    add_footer(s, src)
    return s


def slide_price(prs, bodies, idx, total, out_dir):
    s = blank_slide(prs)
    add_title_band(s, "가격 / 추세", idx, total)
    src = " / ".join(filter(None, (extract_sources(b) for _, b in bodies)))
    chart = build_chart(bodies, out_dir)
    if chart:
        pic_w = Inches(8.8)
        s.shapes.add_picture(chart, MARGIN, CONTENT_TOP, width=pic_w)
        # 우측에 핵심 수치 불릿
        key_lines = []
        for _, body in bodies:
            for ln in clean_bullets(body):
                if any(k in ln for k in ("종가", "MA", "52주", "변동성", "정배열", "거래량")):
                    key_lines.append(ln)
        _, tf = add_textbox(s, MARGIN + pic_w + Inches(0.2), CONTENT_TOP,
                            CONTENT_W - pic_w - Inches(0.2), Inches(4.5))
        for i, ln in enumerate(key_lines[:5]):
            ln = ln if len(ln) <= 34 else ln[:32] + "…"
            write_para(tf, ln, 11, color=INK, first=(i == 0), bullet=True, space_after=8)
        try:
            os.remove(chart)
        except OSError:
            pass
    else:
        # 차트 불가 → 표/불릿 폴백
        tables = []
        for _, body in bodies:
            tables.extend(parse_tables(body))
        if tables:
            render_table(s, tables[0], CONTENT_TOP)
        else:
            return slide_text(prs, "가격 / 추세", bodies, idx, total)
    add_footer(s, src)
    return s


# ---------------------------------------------------------------------------
# PDF 렌더러 (matplotlib 직접 렌더 — 오피스/뷰어 설치 불필요, 브라우저로 열림)
#   PPTX와 동일한 7장 구성·디자인(KB 옐로우/맑은 고딕)을 그대로 재현한다.
# ---------------------------------------------------------------------------
PDF_W, PDF_H = 13.333, 7.5          # inch, 16:9
PX0, PX1 = 0.0525, 0.9475           # 좌/우 본문 경계(figure fraction)
PCW = PX1 - PX0                     # 본문 폭
Y_TITLE, Y_RULE, Y_TOP, Y_FOOT = 0.885, 0.845, 0.795, 0.045
H_YEL = "#FFBC00"
H_INK = "#333333"
H_SUB = "#707070"
H_HEAD = "#333333"
H_ALT = "#F4F4F4"
H_LINE = "#D9D9D9"


def _fp(size, bold=False):
    from matplotlib.font_manager import FontProperties
    fp = FontProperties(fname=MALGUN_TTF) if os.path.exists(MALGUN_TTF) else FontProperties()
    fp.set_size(size)
    if bold:
        fp.set_weight("bold")
    return fp


def _wrap(text, width):
    import textwrap
    return textwrap.wrap(text, width=width) or [""]


def _overlay(fig):
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.axis("off")
    ax.set_zorder(10)
    ax.patch.set_alpha(0)
    return ax


def _rect(ax, x, y, w, h, color):
    import matplotlib.patches as mp
    ax.add_patch(mp.Rectangle((x, y), w, h, color=color, lw=0, zorder=11))


def _pdf_title(ax, title, idx, total):
    _rect(ax, PX0, Y_TITLE - 0.005, 0.009, 0.06, H_YEL)          # 좌측 옐로우 악센트
    ax.text(PX0 + 0.022, Y_TITLE + 0.02, title, color=H_INK,
            fontproperties=_fp(23, bold=True), va="center", ha="left")
    _rect(ax, PX0, Y_RULE, PCW, 0.004, H_YEL)                     # 옐로우 언더라인
    ax.text(PX1, Y_TITLE + 0.02, f"{idx} / {total}", color=H_SUB,
            fontproperties=_fp(11), va="center", ha="right")


def _pdf_footer(ax, source):
    note = "본 자료는 학습용 분석이며 투자 자문이 아닙니다."
    txt = f"{source}   ·   {note}" if source else note
    if len(txt) > 130:
        txt = txt[:128] + "…"
    ax.text(PX0, Y_FOOT, txt, color=H_SUB, fontproperties=_fp(8.5),
            va="center", ha="left")


def _pdf_bullets(ax, lines, size=14, bullet=True, top=Y_TOP, max_lines=MAX_BULLETS):
    line_h = (size / 72) / PDF_H * 1.5
    para_gap = line_h * 0.45
    y = top
    for ln in lines[:max_lines]:
        ln = ln if len(ln) <= 110 else ln[:108] + "…"
        wrapped = _wrap(ln, 48 if bullet else 56)
        if bullet:
            ax.text(PX0, y, "•", color=H_YEL, fontproperties=_fp(size, bold=True),
                    va="top", ha="left")
        for j, wl in enumerate(wrapped):
            ax.text(PX0 + (0.018 if bullet else 0), y, wl, color=H_INK,
                    fontproperties=_fp(size), va="top", ha="left")
            y -= line_h
        y -= para_gap
        if y < Y_FOOT + 0.05:
            break


def _pdf_table(fig, rows, top=Y_TOP, max_rows=MAX_TABLE_ROWS):
    if not rows:
        return top
    header, data = rows[0], rows[1:]
    if len(data) > max_rows:
        data = data[:max_rows] + [["…"] + [""] * (len(header) - 1)]
    body = [header] + data
    n_rows, n_cols = len(body), len(header)
    fsize = 12 if n_cols <= 4 else (10.5 if n_cols <= 6 else 9)
    h = min(0.062 * n_rows, 0.5)
    y0 = top - h
    ax = fig.add_axes([PX0, y0, PCW, h])
    ax.axis("off")
    col_w = [0.30] + [0.70 / (n_cols - 1)] * (n_cols - 1) if n_cols > 1 else [1]
    cells = [[(c or "").replace("**", "") for c in (r + [""] * n_cols)[:n_cols]] for r in body]
    tbl = ax.table(cellText=cells, cellLoc="center", colWidths=col_w, bbox=[0, 0, 1, 1])
    tbl.auto_set_font_size(False)
    tbl.set_fontsize(fsize)
    for (r, c), cell in tbl.get_celld().items():
        cell.set_edgecolor(H_LINE)
        cell.set_linewidth(0.6)
        t = cell.get_text()
        t.set_fontproperties(_fp(fsize, bold=(r == 0 or c == 0)))
        if r == 0:
            cell.set_facecolor(H_HEAD)
            t.set_color("white")
        else:
            cell.set_facecolor(H_ALT if r % 2 == 0 else "white")
            t.set_color(H_INK)
        if c == 0:
            cell.set_text_props(ha="left")
            cell.PAD = 0.04
    return y0 - 0.03


def build_pdf(in_path: str, out_path: str | None = None) -> str:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    from matplotlib.backends.backend_pdf import PdfPages

    with open(in_path, encoding="utf-8") as f:
        md = f.read()
    h1, sections = split_sections(md)
    name = h1 or os.path.splitext(os.path.basename(in_path))[0]
    date, perspective = find_meta(sections)
    if not date:
        date = _dt.date.today().isoformat()
    buckets = collect(sections)
    out_dir = os.path.dirname(os.path.abspath(in_path))
    if out_path is None:
        out_path = os.path.splitext(in_path)[0] + ".pdf"

    def src_of(key):
        return " / ".join(filter(None, (extract_sources(b) for _, b in buckets.get(key, []))))

    def lines_of(key, fallback):
        out = []
        for _, body in buckets.get(key, []):
            out.extend(clean_bullets(body))
        return out or [fallback]

    total = 7
    with PdfPages(out_path) as pdf:
        # 1. 표지
        fig = plt.figure(figsize=(PDF_W, PDF_H))
        ax = _overlay(fig)
        _rect(ax, 0, 0, 0.034, 1, H_YEL)
        _rect(ax, PX0, 0.66, 0.12, 0.006, H_YEL)
        ax.text(PX0, 0.60, name, color=H_INK, fontproperties=_fp(42, bold=True),
                va="center", ha="left")
        ax.text(PX0, 0.52, "종목 리서치 리포트", color=H_SUB,
                fontproperties=_fp(19), va="center", ha="left")
        ax.text(PX0, 0.30, f"작성일 / 기준일 : {date}", color=H_INK,
                fontproperties=_fp(14), va="center", ha="left")
        if perspective:
            ax.text(PX0, 0.24, f"작성 관점 : {perspective}", color=H_SUB,
                    fontproperties=_fp(14), va="center", ha="left")
        _pdf_footer(ax, "")
        pdf.savefig(fig); plt.close(fig)

        # 2. 종목 개요
        fig = plt.figure(figsize=(PDF_W, PDF_H)); ax = _overlay(fig)
        _pdf_title(ax, "종목 개요", 2, total)
        _pdf_bullets(ax, lines_of("overview", "개요 정보가 제공되지 않았습니다."))
        _pdf_footer(ax, src_of("overview")); pdf.savefig(fig); plt.close(fig)

        # 3. 재무 요약 (표 + 코멘트)
        fig = plt.figure(figsize=(PDF_W, PDF_H)); ax = _overlay(fig)
        _pdf_title(ax, "재무 요약 (최근 3개년)", 3, total)
        fin = buckets.get("finance", [("재무", "")])
        tables = []
        for _, body in fin:
            tables.extend(parse_tables(body))
        ynext = Y_TOP
        if tables:
            tables.sort(key=lambda t: (len(t), len(t[0])), reverse=True)
            ynext = _pdf_table(fig, tables[0], Y_TOP)
        comments = [ln for _, b in fin for ln in clean_bullets(b)
                    if "코멘트" in ln or len(ln) > 25]
        if comments:
            _pdf_bullets(ax, comments, size=12.5, top=ynext, max_lines=2)
        _pdf_footer(ax, src_of("finance")); pdf.savefig(fig); plt.close(fig)

        # 4. 가격 / 추세 (차트 + 핵심 수치)
        fig = plt.figure(figsize=(PDF_W, PDF_H)); ax = _overlay(fig)
        _pdf_title(ax, "가격 / 추세", 4, total)
        price = buckets.get("price", [("가격/추세", "")])
        chart = build_chart(price, out_dir)
        if chart:
            img_ax = fig.add_axes([PX0, 0.20, 0.60, 0.58])
            img_ax.imshow(plt.imread(chart)); img_ax.axis("off")
            key = [ln for _, b in price for ln in clean_bullets(b)
                   if any(k in ln for k in ("종가", "MA", "52주", "변동성", "정배열", "거래량"))]
            kx = 0.665
            y = Y_TOP
            for ln in key[:4]:
                ln = re.sub(r"\s*\(출처[^)]*\)", "", ln).strip()  # 우측 요약은 출처 생략(하단 풋터에 표기)
                wrapped = _wrap(ln, 17)[:3]
                ax.text(kx, y, "•", color=H_YEL, fontproperties=_fp(11, bold=True),
                        va="top", ha="left")
                for wl in wrapped:
                    ax.text(kx + 0.016, y, wl, color=H_INK, fontproperties=_fp(11),
                            va="top", ha="left")
                    y -= 0.040
                y -= 0.022
            try:
                os.remove(chart)
            except OSError:
                pass
        else:
            if tables:
                _pdf_table(fig, parse_tables(price[0][1])[0] if parse_tables(price[0][1]) else [], Y_TOP)
            else:
                _pdf_bullets(ax, lines_of("price", "가격 데이터가 제공되지 않았습니다."))
        _pdf_footer(ax, src_of("price")); pdf.savefig(fig); plt.close(fig)

        # 5. 뉴스 · 심리
        fig = plt.figure(figsize=(PDF_W, PDF_H)); ax = _overlay(fig)
        _pdf_title(ax, "뉴스 · 심리", 5, total)
        _pdf_bullets(ax, lines_of("news", "뉴스 정보가 제공되지 않았습니다."))
        _pdf_footer(ax, src_of("news")); pdf.savefig(fig); plt.close(fig)

        # 6. 리스크
        fig = plt.figure(figsize=(PDF_W, PDF_H)); ax = _overlay(fig)
        _pdf_title(ax, "리스크", 6, total)
        _pdf_bullets(ax, lines_of("risk", "리스크 정보가 제공되지 않았습니다."))
        _pdf_footer(ax, src_of("risk")); pdf.savefig(fig); plt.close(fig)

        # 7. 한 줄 종합
        fig = plt.figure(figsize=(PDF_W, PDF_H)); ax = _overlay(fig)
        _pdf_title(ax, "한 줄 종합", 7, total)
        _pdf_bullets(ax, lines_of("summary", "종합 의견이 제공되지 않았습니다."),
                     bullet=False)
        _pdf_footer(ax, src_of("summary")); pdf.savefig(fig); plt.close(fig)

    return out_path


# ---------------------------------------------------------------------------
# 메인
# ---------------------------------------------------------------------------
def build(in_path: str, out_path: str | None = None) -> str:
    with open(in_path, encoding="utf-8") as f:
        md = f.read()

    h1, sections = split_sections(md)
    name = h1 or os.path.splitext(os.path.basename(in_path))[0]
    date, perspective = find_meta(sections)
    if not date:
        date = _dt.date.today().isoformat()
    buckets = collect(sections)

    out_dir = os.path.dirname(os.path.abspath(in_path))
    if out_path is None:
        out_path = os.path.splitext(in_path)[0] + ".pptx"

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 고정 순서
    total = 7
    slide_cover(prs, name, date, perspective)
    slide_text(prs, "종목 개요", buckets.get("overview", [("종목 개요", "개요 정보가 제공되지 않았습니다.")]), 2, total)
    slide_finance(prs, buckets.get("finance", [("재무", "")]), 3, total)
    slide_price(prs, buckets.get("price", [("가격/추세", "")]), 4, total, out_dir)
    slide_text(prs, "뉴스 · 심리", buckets.get("news", [("뉴스·심리", "")]), 5, total)
    slide_text(prs, "리스크", buckets.get("risk", [("리스크", "")]), 6, total)
    slide_text(prs, "한 줄 종합", buckets.get("summary", [("종합", "")]), 7, total, as_bullets=False)

    prs.save(out_path)
    return out_path


if __name__ == "__main__":
    args = [a for a in sys.argv[1:] if not a.startswith("-")]
    flags = [a for a in sys.argv[1:] if a.startswith("-")]
    if not args:
        print("사용법: python build_pptx.py <input.md> [output] [--pdf | --format pptx|pdf|both]")
        sys.exit(1)
    src = args[0]
    dst = args[1] if len(args) > 1 else None

    fmt = "pptx"
    if any(f in ("--pdf",) for f in flags):
        fmt = "both"
    for f in flags:
        if f.startswith("--format"):
            fmt = f.split("=", 1)[1] if "=" in f else "both"
    # `--format pdf` 형태(공백 구분) 지원
    if "--format" in sys.argv:
        i = sys.argv.index("--format")
        if i + 1 < len(sys.argv):
            fmt = sys.argv[i + 1]

    if fmt in ("pptx", "both"):
        print(f"생성 완료: {build(src, dst if (dst and dst.endswith('.pptx')) else None)}")
    if fmt in ("pdf", "both"):
        print(f"생성 완료: {build_pdf(src, dst if (dst and dst.endswith('.pdf')) else None)}")
